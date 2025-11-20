"""
Generate editable PowerPoint figures for AAAI paper
All shapes, arrows, and text are fully editable
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_figure_2():
    """Figure 2: System Architecture"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title.text_frame
    tf.text = "System Architecture: RL-Enhanced MPC Framework"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Box 1: RL Optimizer
    box1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.2), Inches(2.8), Inches(2.2)
    )
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(227, 242, 253)  # Light blue
    box1.line.color.rgb = RGBColor(25, 118, 210)
    box1.line.width = Pt(3)

    text_frame = box1.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "RL Optimizer (PPO)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(13, 71, 161)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "\nState Space (29D):"
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "• Tracking errors"
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "• Control effort"
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "• Current hyperparameters"
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "\nAction: Q, R, N (17D)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Box 2: MPC Controller
    box2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.6), Inches(1.2), Inches(2.8), Inches(2.2)
    )
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(232, 245, 233)  # Light green
    box2.line.color.rgb = RGBColor(56, 142, 60)
    box2.line.width = Pt(3)

    text_frame = box2.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "MPC Controller"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(27, 94, 32)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "(CasADi/IPOPT)"
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "\nOptimization:"
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "• Cost function J(Q,R,N)"
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "• Dynamics constraints"
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "• Control limits"
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

    # Box 3: UAV Environment
    box3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.7), Inches(1.2), Inches(2.8), Inches(2.2)
    )
    box3.fill.solid()
    box3.fill.fore_color.rgb = RGBColor(255, 243, 224)  # Light orange
    box3.line.color.rgb = RGBColor(245, 124, 0)
    box3.line.width = Pt(3)

    text_frame = box3.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "UAV Environment"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(230, 81, 0)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "(PyBullet)"
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "\nState Space (12D):"
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "• Position [x, y, z]"
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "• Velocity [vx, vy, vz]"
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "• Angles & Rates"
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

    # Arrow 1: RL to MPC
    arrow1 = slide.shapes.add_connector(1, Inches(3.3), Inches(2.2), Inches(3.6), Inches(2.2))
    arrow1.line.color.rgb = RGBColor(25, 118, 210)
    arrow1.line.width = Pt(3)

    label1 = slide.shapes.add_textbox(Inches(3.3), Inches(1.8), Inches(0.3), Inches(0.3))
    tf = label1.text_frame
    tf.text = "Q, R, N"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(25, 118, 210)

    # Arrow 2: MPC to UAV
    arrow2 = slide.shapes.add_connector(1, Inches(6.4), Inches(2.2), Inches(6.7), Inches(2.2))
    arrow2.line.color.rgb = RGBColor(25, 118, 210)
    arrow2.line.width = Pt(3)

    label2 = slide.shapes.add_textbox(Inches(6.4), Inches(1.8), Inches(0.3), Inches(0.3))
    tf = label2.text_frame
    tf.text = "u(t)"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(25, 118, 210)

    # Arrow 3: UAV to MPC (feedback)
    arrow3 = slide.shapes.add_connector(1, Inches(6.7), Inches(3.0), Inches(6.4), Inches(3.0))
    arrow3.line.color.rgb = RGBColor(211, 47, 47)
    arrow3.line.width = Pt(3)

    label3 = slide.shapes.add_textbox(Inches(6.4), Inches(3.2), Inches(0.3), Inches(0.3))
    tf = label3.text_frame
    tf.text = "x(t)"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(211, 47, 47)

    # Arrow 4: UAV to RL (bottom feedback)
    arrow4 = slide.shapes.add_connector(1, Inches(8.1), Inches(3.4), Inches(8.1), Inches(4.5))
    arrow4.line.color.rgb = RGBColor(211, 47, 47)
    arrow4.line.width = Pt(3)

    arrow5 = slide.shapes.add_connector(1, Inches(8.1), Inches(4.5), Inches(1.9), Inches(4.5))
    arrow5.line.color.rgb = RGBColor(211, 47, 47)
    arrow5.line.width = Pt(3)

    arrow6 = slide.shapes.add_connector(1, Inches(1.9), Inches(4.5), Inches(1.9), Inches(3.4))
    arrow6.line.color.rgb = RGBColor(211, 47, 47)
    arrow6.line.width = Pt(3)

    label4 = slide.shapes.add_textbox(Inches(3.5), Inches(4.6), Inches(3), Inches(0.3))
    tf = label4.text_frame
    tf.text = "Performance Metrics"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(211, 47, 47)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Notes
    note1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(5.2), Inches(2.8), Inches(1.2)
    )
    note1.fill.solid()
    note1.fill.fore_color.rgb = RGBColor(245, 245, 245)
    note1.line.color.rgb = RGBColor(25, 118, 210)
    note1.line.width = Pt(2)
    note1.line.dash_style = 2  # Dashed

    text_frame = note1.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "PPO Algorithm"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 118, 210)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "• 2 hidden layers (256 units)"
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "• 4 parallel environments"
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "• Learning rate: 3×10⁻⁴"
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

    note2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.6), Inches(5.2), Inches(2.8), Inches(1.2)
    )
    note2.fill.solid()
    note2.fill.fore_color.rgb = RGBColor(245, 245, 245)
    note2.line.color.rgb = RGBColor(56, 142, 60)
    note2.line.width = Pt(2)
    note2.line.dash_style = 2

    text_frame = note2.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "Real-Time MPC"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(56, 142, 60)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "• Solve time: 30-40ms"
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "• Horizon N: 10 steps"
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "• Update rate: 50Hz"
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

    note3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.7), Inches(5.2), Inches(2.8), Inches(1.2)
    )
    note3.fill.solid()
    note3.fill.fore_color.rgb = RGBColor(245, 245, 245)
    note3.line.color.rgb = RGBColor(245, 124, 0)
    note3.line.width = Pt(2)
    note3.line.dash_style = 2

    text_frame = note3.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "High-Fidelity Simulation"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(245, 124, 0)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "• Physics Δt: 1ms"
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "• Control Δt: 20ms"
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "• RK4 integration"
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

    prs.save('D:/rl_tuned_mpc/paper/figure_2_architecture.pptx')
    print("[OK] Created: figure_2_architecture.pptx")


def create_figure_3():
    """Figure 3: Transfer Learning Flow"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(11)  # Taller slide for vertical layout

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title.text_frame
    tf.text = "Sequential Transfer Learning Across UAV Platforms"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Phase 1: Crazyflie
    box1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1), Inches(8), Inches(1.5)
    )
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(187, 222, 251)
    box1.line.color.rgb = RGBColor(25, 118, 210)
    box1.line.width = Pt(4)

    text_frame = box1.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "Phase 1: Crazyflie 2.X (0.027 kg) - Base Training"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "Steps: 20,000 | Learning rate: 3×10⁻⁴ | Time: 200 min | From scratch"
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

    # Arrow
    arr1 = slide.shapes.add_connector(1, Inches(5), Inches(2.5), Inches(5), Inches(3))
    arr1.line.color.rgb = RGBColor(25, 118, 210)
    arr1.line.width = Pt(4)

    # Phase 2: Racing
    box2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(3), Inches(8), Inches(1.5)
    )
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(144, 202, 249)
    box2.line.color.rgb = RGBColor(25, 118, 210)
    box2.line.width = Pt(4)

    text_frame = box2.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "Phase 2: Racing Drone (0.800 kg) - Fine-Tuning"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "Steps: 5,000 (25%) | Learning rate: 3×10⁻⁵ | Time: 52 min | From Crazyflie"
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

    # Arrow
    arr2 = slide.shapes.add_connector(1, Inches(5), Inches(4.5), Inches(5), Inches(5))
    arr2.line.color.rgb = RGBColor(25, 118, 210)
    arr2.line.width = Pt(4)

    # Phase 3: Generic
    box3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(5), Inches(8), Inches(1.5)
    )
    box3.fill.solid()
    box3.fill.fore_color.rgb = RGBColor(100, 181, 246)
    box3.line.color.rgb = RGBColor(25, 118, 210)
    box3.line.width = Pt(4)

    text_frame = box3.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "Phase 3: Generic Quadrotor (2.500 kg) - Fine-Tuning"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "Steps: 5,000 (25%) | Learning rate: 3×10⁻⁵ | Time: 52 min | From Racing"
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

    # Arrow
    arr3 = slide.shapes.add_connector(1, Inches(5), Inches(6.5), Inches(5), Inches(7))
    arr3.line.color.rgb = RGBColor(25, 118, 210)
    arr3.line.width = Pt(4)

    # Phase 4: Heavy-Lift
    box4 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(7), Inches(8), Inches(1.5)
    )
    box4.fill.solid()
    box4.fill.fore_color.rgb = RGBColor(66, 165, 245)
    box4.line.color.rgb = RGBColor(25, 118, 210)
    box4.line.width = Pt(4)

    text_frame = box4.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "Phase 4: Heavy-Lift Hexarotor (5.500 kg) - Fine-Tuning"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "Steps: 5,000 (25%) | Learning rate: 3×10⁻⁵ | Time: 59 min | From Generic"
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

    # Arrow to summary
    arr4 = slide.shapes.add_connector(1, Inches(5), Inches(8.5), Inches(5), Inches(9))
    arr4.line.color.rgb = RGBColor(56, 142, 60)
    arr4.line.width = Pt(4)

    # Summary box
    summary = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(9), Inches(8), Inches(1.7)
    )
    summary.fill.solid()
    summary.fill.fore_color.rgb = RGBColor(200, 230, 201)
    summary.line.color.rgb = RGBColor(56, 142, 60)
    summary.line.width = Pt(4)

    text_frame = summary.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "Summary: Efficiency Gains"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(27, 94, 32)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "Without Transfer: 80,000 steps, 828 min | With Transfer: 35,000 steps, 363 min"
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "✓ 75% step reduction  ✓ 56.2% time savings  ✓ 1.34±0.01m RMSE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    prs.save('D:/rl_tuned_mpc/paper/figure_3_transfer.pptx')
    print("[OK] Created: figure_3_transfer.pptx")


def create_figure_4():
    """Figure 4: PyBullet Environment"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title.text_frame
    tf.text = "PyBullet UAV Simulation Environment"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Control Input
    box1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.2), Inches(2.2), Inches(1.5)
    )
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(232, 245, 233)
    box1.line.color.rgb = RGBColor(56, 142, 60)
    box1.line.width = Pt(3)

    text_frame = box1.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "Control Input u(t)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "4D: Thrust, Roll/\nPitch/Yaw rates"
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

    # UAV
    box2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.2), Inches(1.2), Inches(3.6), Inches(1.5)
    )
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(227, 242, 253)
    box2.line.color.rgb = RGBColor(25, 118, 210)
    box2.line.width = Pt(3)

    text_frame = box2.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "Quadrotor UAV"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "4 Motors + Propellers | Body: Mass & Inertia | 6-DOF"
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

    # Physics
    box3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.3), Inches(1.2), Inches(2.2), Inches(1.5)
    )
    box3.fill.solid()
    box3.fill.fore_color.rgb = RGBColor(225, 245, 254)
    box3.line.color.rgb = RGBColor(2, 136, 209)
    box3.line.width = Pt(3)

    text_frame = box3.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "PyBullet Physics"
    p.font.size = Pt(13)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "RK4, Δt=1ms\n240Hz, Rigid body"
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

    # Environment
    box4 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.3), Inches(3.2), Inches(2.2), Inches(1.3)
    )
    box4.fill.solid()
    box4.fill.fore_color.rgb = RGBColor(243, 229, 245)
    box4.line.color.rgb = RGBColor(123, 31, 162)
    box4.line.width = Pt(3)

    text_frame = box4.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "Environment"
    p.font.size = Pt(13)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "Aero drag, Ground\neffect, Collision"
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

    # State
    box5 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.2), Inches(3.2), Inches(3.6), Inches(1.5)
    )
    box5.fill.solid()
    box5.fill.fore_color.rgb = RGBColor(255, 243, 224)
    box5.line.color.rgb = RGBColor(245, 124, 0)
    box5.line.width = Pt(3)

    text_frame = box5.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "State Vector x(t)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "12D: Position, Velocity, Angles, Rates"
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

    # Feedback
    box6 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(3.2), Inches(2.2), Inches(1.2)
    )
    box6.fill.solid()
    box6.fill.fore_color.rgb = RGBColor(255, 249, 196)
    box6.line.color.rgb = RGBColor(249, 168, 37)
    box6.line.width = Pt(3)

    text_frame = box6.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "Feedback Loop"
    p.font.size = Pt(13)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "Error e(t)\nTo MPC"
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

    # Arrows (simplified - you can adjust in PowerPoint)
    arr1 = slide.shapes.add_connector(1, Inches(2.7), Inches(1.9), Inches(3.2), Inches(1.9))
    arr1.line.color.rgb = RGBColor(25, 118, 210)
    arr1.line.width = Pt(3)

    arr2 = slide.shapes.add_connector(1, Inches(6.8), Inches(1.9), Inches(7.3), Inches(1.9))
    arr2.line.color.rgb = RGBColor(25, 118, 210)
    arr2.line.width = Pt(3)

    arr3 = slide.shapes.add_connector(1, Inches(8.4), Inches(2.7), Inches(8.4), Inches(3.2))
    arr3.line.color.rgb = RGBColor(25, 118, 210)
    arr3.line.width = Pt(3)

    arr4 = slide.shapes.add_connector(1, Inches(7.3), Inches(3.8), Inches(6.8), Inches(3.8))
    arr4.line.color.rgb = RGBColor(25, 118, 210)
    arr4.line.width = Pt(3)

    arr5 = slide.shapes.add_connector(1, Inches(3.2), Inches(3.8), Inches(2.7), Inches(3.8))
    arr5.line.color.rgb = RGBColor(25, 118, 210)
    arr5.line.width = Pt(3)

    arr6 = slide.shapes.add_connector(1, Inches(1.6), Inches(3.2), Inches(1.6), Inches(2.7))
    arr6.line.color.rgb = RGBColor(211, 47, 47)
    arr6.line.width = Pt(3)

    # Info boxes
    info1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(5.5), Inches(3), Inches(1)
    )
    info1.fill.solid()
    info1.fill.fore_color.rgb = RGBColor(250, 250, 250)
    info1.line.color.rgb = RGBColor(25, 118, 210)
    info1.line.width = Pt(2)

    text_frame = info1.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "Platform Parameters"
    p.font.size = Pt(11)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "Mass: 0.027-5.5 kg (200×)\nInertia: 6000× variation"
    p.font.size = Pt(9)
    p.alignment = PP_ALIGN.CENTER

    info2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.5), Inches(5.5), Inches(3), Inches(1)
    )
    info2.fill.solid()
    info2.fill.fore_color.rgb = RGBColor(250, 250, 250)
    info2.line.color.rgb = RGBColor(2, 136, 209)
    info2.line.width = Pt(2)

    text_frame = info2.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "Simulation Fidelity"
    p.font.size = Pt(11)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "High-fidelity 3D simulation\nSim-to-real ready"
    p.font.size = Pt(9)
    p.alignment = PP_ALIGN.CENTER

    info3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.5), Inches(5.5), Inches(3), Inches(1)
    )
    info3.fill.solid()
    info3.fill.fore_color.rgb = RGBColor(250, 250, 250)
    info3.line.color.rgb = RGBColor(123, 31, 162)
    info3.line.width = Pt(2)

    text_frame = info3.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "Realistic Physics"
    p.font.size = Pt(11)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "Validated dynamics\nGround interaction"
    p.font.size = Pt(9)
    p.alignment = PP_ALIGN.CENTER

    prs.save('D:/rl_tuned_mpc/paper/figure_4_pybullet.pptx')
    print("[OK] Created: figure_4_pybullet.pptx")


if __name__ == '__main__':
    print("="*70)
    print("Generating Editable PowerPoint Figures")
    print("="*70)
    print()

    create_figure_2()
    create_figure_3()
    create_figure_4()

    print()
    print("="*70)
    print("[SUCCESS] All 3 PowerPoint files created!")
    print("="*70)
    print()
    print("Files created:")
    print("  1. figure_2_architecture.pptx")
    print("  2. figure_3_transfer.pptx")
    print("  3. figure_4_pybullet.pptx")
    print()
    print("Location: D:\\rl_tuned_mpc\\paper\\")
    print()
    print("Next steps:")
    print("  1. Open each .pptx file in PowerPoint")
    print("  2. Edit shapes, arrows, text as needed")
    print("  3. Export as PDF:")
    print("     File → Save As → PDF")
    print("  4. Or export as PNG:")
    print("     File → Save As → PNG")
    print()
